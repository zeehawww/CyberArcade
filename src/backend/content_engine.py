"""
Content Engine: file parsing, feed fetching, and AI content generation.
"""
import re, json, random, textwrap
from datetime import datetime

# ── File Parsers ──────────────────────────────────────────────
def extract_text_from_pdf(file_bytes):
    """Extract text from PDF bytes using PyPDF2 or pdfplumber."""
    try:
        import pdfplumber, io
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            return '\n'.join(p.extract_text() or '' for p in pdf.pages)
    except ImportError:
        pass
    try:
        from PyPDF2 import PdfReader
        import io
        reader = PdfReader(io.BytesIO(file_bytes))
        return '\n'.join(p.extract_text() or '' for p in reader.pages)
    except ImportError:
        return '[PDF parsing libraries not installed. Install pdfplumber or PyPDF2]'

def extract_text_from_docx(file_bytes):
    try:
        from docx import Document
        import io
        doc = Document(io.BytesIO(file_bytes))
        return '\n'.join(p.text for p in doc.paragraphs)
    except ImportError:
        return '[python-docx not installed. Run: pip install python-docx]'

def extract_text_from_pptx(file_bytes):
    try:
        from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        return '\n'.join(texts)
    except ImportError:
        return '[python-pptx not installed. Run: pip install python-pptx]'

# ── Feed Fetchers ─────────────────────────────────────────────
def fetch_feed(url, feed_type):
    """Fetch items from an external feed. Returns list of {title, body, link}."""
    import urllib.request
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'CyberArcade/1.0'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = resp.read()
    except Exception as e:
        return [{'title': 'Feed Error', 'body': str(e), 'link': url}]

    if feed_type == 'nist':
        return _parse_nist_json(data)
    else:
        return _parse_rss(data)

def _parse_rss(data):
    import xml.etree.ElementTree as ET
    items = []
    try:
        root = ET.fromstring(data)
        for item in root.iter('item'):
            t = item.findtext('title', '')
            d = item.findtext('description', '')
            l = item.findtext('link', '')
            items.append({'title': t, 'body': _strip_html(d)[:2000], 'link': l})
        if not items:
            for entry in root.iter('{http://www.w3.org/2005/Atom}entry'):
                t = entry.findtext('{http://www.w3.org/2005/Atom}title', '')
                d = entry.findtext('{http://www.w3.org/2005/Atom}summary', '')
                lnk = entry.find('{http://www.w3.org/2005/Atom}link')
                l = lnk.get('href','') if lnk is not None else ''
                items.append({'title': t, 'body': _strip_html(d)[:2000], 'link': l})
    except Exception:
        pass
    return items[:10]

def _parse_nist_json(data):
    items = []
    try:
        obj = json.loads(data)
        for vuln in obj.get('vulnerabilities', [])[:10]:
            cve = vuln.get('cve', {})
            cid = cve.get('id', '')
            descs = cve.get('descriptions', [])
            desc = next((d['value'] for d in descs if d.get('lang') == 'en'), '')
            items.append({'title': cid, 'body': desc[:2000], 'link': f'https://nvd.nist.gov/vuln/detail/{cid}'})
    except Exception:
        pass
    return items

def _strip_html(text):
    return re.sub(r'<[^>]+>', '', text or '')

# ── Content Generation Engine ─────────────────────────────────
# These generators work locally without requiring external AI APIs.
# They use rule-based templates for deterministic, instant output.

CATEGORIES = ['phishing','ransomware','social_engineering','password_security','data_privacy','zero_day_threats']

def categorize_content(text):
    """Auto-categorize content into one of the 6 threat categories."""
    text_lower = text.lower()
    scores = {}
    kw = {
        'phishing': ['phishing','spear','email','bec','credential','harvest','lure','click'],
        'ransomware': ['ransomware','encrypt','ransom','lockbit','conti','decrypt','extort'],
        'social_engineering': ['social engineer','pretext','vishing','tailgat','impersonat','manipulat'],
        'password_security': ['password','credential','brute','mfa','2fa','authenticat','token'],
        'data_privacy': ['privacy','gdpr','pii','data breach','compliance','hipaa','leak'],
        'zero_day_threats': ['zero.day','0-day','exploit','vulnerability','cve','patch','buffer'],
    }
    for cat, words in kw.items():
        scores[cat] = sum(1 for w in words if w in text_lower)
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else 'general'

def simplify_for_tone(text, tone):
    """Rewrite content summary for a specific audience tone."""
    short = text[:500]
    if tone == 'child':
        return f"🛡️ Hey kids! Here's something important about staying safe online:\n\n{_childify(short)}\n\n💡 Remember: Always ask a grown-up if something online feels weird!"
    elif tone == 'college':
        return f"📚 Academic Analysis:\n\n{short}\n\nKey Takeaway: This topic intersects with information security principles including the CIA triad. Consider how institutional policies address this threat vector."
    elif tone == 'business':
        return f"📊 Business Impact Brief:\n\n{short}\n\n⚠️ Risk Level: Medium-High\n💰 Potential Impact: Operational disruption, regulatory fines\n✅ Action Required: Review organizational policies and employee training."
    elif tone == 'technical':
        return f"🔧 Technical Analysis:\n\n{short}\n\n🔍 IOCs to monitor: Check SIEM logs for anomalous patterns.\n🛡️ Mitigations: Apply defense-in-depth, update WAF rules, enable enhanced logging."
    return short

def _childify(text):
    replacements = [
        ('vulnerability','weak spot'),('exploit','trick'),('malware','bad software'),
        ('credential','password'),('authentication','login check'),('encryption','secret code'),
        ('unauthorized','not allowed'),('compromise','break into'),('exfiltrate','steal'),
    ]
    result = text
    for old, new in replacements:
        result = re.sub(old, new, result, flags=re.IGNORECASE)
    return result

def generate_micro_learning(title, text):
    cat = categorize_content(text)
    short = text[:300]
    return {
        'type': 'micro_learning',
        'title': f'⚡ Quick Learn: {title}',
        'body': json.dumps({
            'duration': '2-3 minutes',
            'category': cat,
            'objective': f'Understand the basics of {title}',
            'key_points': [s.strip()+'.' for s in short.split('.')[:3] if s.strip()],
            'action_item': f'Review your organization\'s policy on {cat.replace("_"," ")}.',
            'quiz': _make_quiz(title, cat)
        }, indent=2)
    }

def generate_infographic(title, text):
    cat = categorize_content(text)
    sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 15]
    # Extract meaningful labels from sentences
    label_prefixes = _infographic_labels(cat)
    stats = []
    for i, s in enumerate(sentences[:5]):
        label = label_prefixes[i % len(label_prefixes)]
        stats.append({'label': label, 'value': s[:100]})
    if not stats:
        stats = [{'label': 'Overview', 'value': text[:100]}]
    return {
        'type': 'infographic',
        'title': f'📊 Infographic: {title}',
        'body': json.dumps({
            'headline': title,
            'category': cat,
            'icon': _cat_icon(cat),
            'stats': stats,
            'color_scheme': _cat_colors(cat),
            'footer': f'Source: CyberArcade Threat Intelligence | {cat.replace("_"," ").title()} Analysis'
        }, indent=2)
    }

def generate_poster(title, text):
    cat = categorize_content(text)
    # Extract the most impactful sentence as key message
    sentences = [s.strip() for s in text.split('.') if 15 < len(s.strip()) < 200]
    key_msg = sentences[0] if sentences else text[:150]
    return {
        'type': 'poster',
        'title': f'🖼️ Poster: {title}',
        'body': json.dumps({
            'headline': title[:60],
            'tagline': random.choice(_poster_taglines(cat)),
            'category': cat,
            'visual_style': 'social_media',
            'dimensions': '1080x1080',
            'key_message': key_msg[:180],
            'call_to_action': random.choice(_poster_ctas(cat)),
            'color_scheme': _cat_colors(cat)
        }, indent=2)
    }

def generate_explainer_script(title, text):
    cat = categorize_content(text)
    short = text[:400]
    return {
        'type': 'explainer_script',
        'title': f'🎬 Explainer: {title}',
        'body': json.dumps({
            'duration': '60-90 seconds',
            'script': {
                'intro': f'[NARRATOR] Today we\'re talking about {title}. This is something everyone should know about.',
                'body': f'[NARRATOR] {short}',
                'conclusion': f'[NARRATOR] Remember - staying safe online is everyone\'s responsibility. Keep learning at CyberArcade!',
            },
            'visual_notes': ['Title card with animated shield icon', f'Animated diagram of {cat.replace("_"," ")}', 'Key points appear as bullet list', 'End card with CyberArcade logo']
        }, indent=2)
    }

def generate_quiz_questions(title, text):
    cat = categorize_content(text)
    return {
        'type': 'quiz',
        'title': f'❓ Quiz: {title}',
        'body': json.dumps({
            'category': cat,
            'questions': _make_quiz_set(title, cat, text)
        }, indent=2)
    }

def generate_scenario_simulation(title, text):
    cat = categorize_content(text)
    return {
        'type': 'scenario',
        'title': f'🎯 Scenario: {title}',
        'body': json.dumps({
            'category': cat,
            'scenario': _make_scenario(cat),
            'difficulty': 'intermediate'
        }, indent=2)
    }

def generate_video_storyboard(title, text):
    cat = categorize_content(text)
    cat_display = cat.replace('_', ' ')
    sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 20]
    main_point = sentences[0] if sentences else text[:200]
    detail = sentences[1] if len(sentences) > 1 else ''
    return {
        'type': 'video_storyboard',
        'title': f'🎥 Video: {title}',
        'body': json.dumps({
            'duration': '2-3 minutes',
            'target_audience': 'All employees',
            'scenes': [
                {'scene': 1, 'visual': f'Motion graphics title sequence with {cat_display} iconography', 'narration': f'In today\'s threat landscape, understanding {title.lower()} is critical for every organization.', 'duration': '15s'},
                {'scene': 2, 'visual': f'Animated attack flow diagram showing {cat_display} methodology', 'narration': main_point[:250], 'duration': '35s'},
                {'scene': 3, 'visual': 'Split-screen: real-world examples vs. warning indicators', 'narration': detail[:200] if detail else f'Recognizing the warning signs of {cat_display} can prevent significant damage to your organization.', 'duration': '30s'},
                {'scene': 4, 'visual': 'Checklist animation with actionable defense steps', 'narration': f'Here are the critical steps your team should implement: verify all requests through trusted channels, maintain updated security controls, and report suspicious activity immediately.', 'duration': '25s'},
                {'scene': 5, 'visual': 'Professional end card with key takeaway and CyberArcade branding', 'narration': f'Remember: {cat_display} threats evolve constantly. Stay informed, stay vigilant, and make security awareness part of your daily routine.', 'duration': '15s'},
            ]
        }, indent=2)
    }

def generate_comic(title, text):
    cat = categorize_content(text)
    return {
        'type': 'comic',
        'title': f'📖 Comic: {title}',
        'body': json.dumps({
            'panels': [
                {'panel': 1, 'character': '🦸 Cyber Hero', 'dialogue': f'Oh no! I see a {cat.replace("_"," ")} attack happening!', 'scene': 'Hero at computer'},
                {'panel': 2, 'character': '🦹 Villain', 'dialogue': _villain_line(cat), 'scene': 'Dark figure behind screen'},
                {'panel': 3, 'character': '🦸 Cyber Hero', 'dialogue': f'Not so fast! I know how to handle {cat.replace("_"," ")}!', 'scene': 'Hero activating shield'},
                {'panel': 4, 'character': '🦸 Cyber Hero', 'dialogue': _hero_tip(cat), 'scene': 'Hero teaching friends'},
                {'panel': 5, 'character': '👨‍👩‍👧‍👦 Everyone', 'dialogue': 'Thanks Cyber Hero! Now we know how to stay safe!', 'scene': 'Happy ending'},
            ],
            'style': 'child_friendly',
            'age_group': '6-12'
        }, indent=2)
    }

def generate_executive_dashboard(title, text):
    cat = categorize_content(text)
    cat_display = cat.replace('_', ' ').title()
    # Derive severity from content keywords
    high_risk_words = ['critical', 'severe', 'urgent', 'immediate', 'breach', 'compromised', 'active exploitation']
    risk_score = sum(1 for w in high_risk_words if w in text.lower())
    risk_level = 'Critical' if risk_score >= 3 else 'High' if risk_score >= 2 else 'Medium' if risk_score >= 1 else 'Low'
    word_count = len(text.split())
    sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 10]
    return {
        'type': 'executive_dashboard',
        'title': f'📈 Executive Summary: {title}',
        'body': json.dumps({
            'threat_category': cat_display,
            'risk_level': risk_level,
            'summary': text[:300],
            'metrics': {
                'content_depth': f'{word_count} words analyzed',
                'key_findings': f'{min(len(sentences), 10)} actionable insights',
                'risk_score': f'{min(risk_score * 25 + random.randint(15, 35), 100)}/100',
                'response_urgency': risk_level,
            },
            'recommendations': [
                f'Conduct targeted {cat_display.lower()} awareness training for all staff',
                f'Review and update {cat_display.lower()} detection and response procedures',
                f'Assess current controls against the specific {cat_display.lower()} vectors identified',
                'Schedule a tabletop exercise simulating this threat scenario',
                'Brief executive leadership on risk exposure and mitigation timeline'
            ],
            'trend': 'increasing' if risk_score >= 2 else 'stable'
        }, indent=2)
    }

# ── Helpers ───────────────────────────────────────────────────
def _cat_icon(cat):
    icons = {'phishing':'🎣','ransomware':'🔒','social_engineering':'🎭','password_security':'🔑','data_privacy':'🛡️','zero_day_threats':'⚡'}
    return icons.get(cat, '🔐')

def _cat_colors(cat):
    colors = {
        'phishing':['#ff6b6b','#ee5a24'],
        'ransomware':['#a55eea','#8854d0'],
        'social_engineering':['#fd9644','#e67e22'],
        'password_security':['#2bcbba','#0fb9b1'],
        'data_privacy':['#4b7bec','#3867d6'],
        'zero_day_threats':['#fc5c65','#eb3b5a'],
    }
    return colors.get(cat, ['#00ffff','#0080ff'])

def _infographic_labels(cat):
    labels = {
        'phishing': ['Attack Vector', 'Threat Indicator', 'Impact Assessment', 'Mitigation Strategy', 'Detection Method'],
        'ransomware': ['Infection Mechanism', 'Encryption Method', 'Recovery Protocol', 'Prevention Measure', 'Incident Timeline'],
        'social_engineering': ['Manipulation Tactic', 'Psychological Trigger', 'Target Profile', 'Defense Protocol', 'Warning Sign'],
        'password_security': ['Vulnerability Factor', 'Authentication Layer', 'Credential Risk', 'Best Practice', 'Compliance Standard'],
        'data_privacy': ['Data Classification', 'Regulatory Requirement', 'Exposure Risk', 'Protection Measure', 'Compliance Gap'],
        'zero_day_threats': ['Exploit Vector', 'Vulnerability Scope', 'Patch Priority', 'Detection Signal', 'Response Protocol'],
    }
    return labels.get(cat, ['Key Finding', 'Analysis', 'Impact', 'Recommendation', 'Action Item'])

def _poster_taglines(cat):
    tags = {
        'phishing': ['Think Before You Click.', 'One Click Can Cost Millions.', 'Verify. Don\'t Trust Blindly.', 'Suspicious Link? Stop. Report. Delete.'],
        'ransomware': ['Backup Today, Recover Tomorrow.', 'Your Data Is Your Responsibility.', 'Don\'t Negotiate With Cyber Criminals.', 'Prepare. Prevent. Protect.'],
        'social_engineering': ['Trust Is Earned, Not Assumed.', 'Verify Every Request. Every Time.', 'The Weakest Link Is Human.', 'Question Everything. Protect Everyone.'],
        'password_security': ['Strong Passwords. Stronger Defense.', 'Your Password Is Your First Firewall.', '123456 Is Not A Password.', 'Multi-Factor: Because One Layer Isn\'t Enough.'],
        'data_privacy': ['Your Data. Your Rights. Your Rules.', 'Privacy Is Not Optional.', 'Protect What Matters Most.', 'Data Breaches Start With Negligence.'],
        'zero_day_threats': ['Patch Now. Thank Yourself Later.', 'Unknown Threats Require Known Defenses.', 'Every Delay Is An Opportunity For Attackers.', 'Vigilance Is The Best Zero-Day Defense.'],
    }
    return tags.get(cat, ['Stay Secure. Stay Informed.', 'Cybersecurity Is Everyone\'s Job.'])

def _poster_ctas(cat):
    ctas = {
        'phishing': ['Report suspicious emails to your IT team immediately', 'Enable email filtering and train your team today', 'Take the CyberArcade phishing awareness course'],
        'ransomware': ['Verify your backup strategy with IT this week', 'Schedule a ransomware readiness assessment', 'Review your incident response plan now'],
        'social_engineering': ['Implement a verification protocol for all sensitive requests', 'Conduct social engineering awareness training', 'Establish a zero-trust verification culture'],
        'password_security': ['Enable MFA on all critical accounts today', 'Audit your password policy against NIST guidelines', 'Deploy a password manager organization-wide'],
        'data_privacy': ['Conduct a data privacy audit this quarter', 'Review your data handling procedures now', 'Ensure GDPR/CCPA compliance across all systems'],
        'zero_day_threats': ['Enable automatic security updates on all endpoints', 'Subscribe to vulnerability alerts from CISA and NIST', 'Implement network segmentation to limit blast radius'],
    }
    return ctas.get(cat, ['Strengthen your security posture with CyberArcade', 'Start your cybersecurity training today'])

def _villain_line(cat):
    lines = {
        'phishing': random.choice([
            'I crafted the perfect email — your CEO\'s name, urgent language, a link that looks legitimate. Will your team fall for it?',
            'One spoofed domain, one convincing pretext. That\'s all it takes to harvest credentials from hundreds of employees.'
        ]),
        'ransomware': random.choice([
            'Your files are now encrypted with military-grade AES-256. The clock is ticking — pay up or lose everything.',
            'I exploited one unpatched server. Now your entire network is locked. How much is your data worth to you?'
        ]),
        'social_engineering': random.choice([
            'I don\'t need to hack your systems. I just need to convince one employee I\'m from HR. Humans are always the weakest link.',
            'A friendly phone call, a fake badge, a confident walk through the door — physical and digital access, no hacking required.'
        ]),
        'password_security': random.choice([
            'Password123? Company name plus the year? I have a dictionary of 10 billion passwords. Yours took 0.3 seconds.',
            'One credential stuffing attack with leaked databases. 65% of your employees reuse passwords across services.'
        ]),
        'data_privacy': random.choice([
            'Your customer database was left on an unencrypted S3 bucket. I\'ve already downloaded 2 million records.',
            'Third-party tracking, unencrypted PII, no consent management — your compliance violations are my goldmine.'
        ]),
        'zero_day_threats': random.choice([
            'I found a vulnerability in your web framework before the vendor did. Zero patch. Zero defense. Full access.',
            'While you wait for the vendor to release a fix, I\'ve already weaponized the exploit and sold it on the dark web.'
        ]),
    }
    return lines.get(cat, 'Your defenses have gaps you don\'t even know about. I\'ve already mapped them all.')

def _hero_tip(cat):
    tips = {
        'phishing': 'Always verify the sender\'s email address — not just the display name. Hover over links to check the actual URL. When in doubt, contact the sender through a separate, trusted channel.',
        'ransomware': 'Maintain offline backups following the 3-2-1 rule: 3 copies, 2 different media types, 1 offsite. Keep all software patched and segment your network to contain potential breaches.',
        'social_engineering': 'Establish verification procedures for sensitive requests. Never share credentials or grant access based on a phone call alone. Always verify through official, independent channels.',
        'password_security': 'Use a password manager and enable multi-factor authentication on every account. Create unique passphrases of 16+ characters. Never reuse passwords across different services.',
        'data_privacy': 'Classify your data by sensitivity level. Encrypt PII at rest and in transit. Review app permissions regularly and practice data minimization — only collect what you truly need.',
        'zero_day_threats': 'Enable automatic updates and subscribe to threat intelligence feeds. Use application whitelisting and network segmentation to reduce your attack surface against unknown exploits.',
    }
    return tips.get(cat, 'Stay informed, stay vigilant, and report anything suspicious to your security team immediately.')

def _make_quiz(title, cat):
    return {'question': f'What is the most effective defense against {cat.replace("_"," ")} threats?',
            'options': ['Ignore suspicious activity and hope for the best', 'Implement layered security controls and continuous training', 'Rely solely on antivirus software', 'Disconnect from the internet entirely'],
            'correct': 1}

def _make_quiz_set(title, cat, text):
    cat_display = cat.replace('_', ' ')
    questions = [
        _make_quiz(title, cat),
        {'question': f'You encounter a potential {cat_display} incident. What is the correct first step?',
         'options': ['Try to fix it yourself', 'Immediately report it to your security team', 'Ignore it if it seems minor', 'Post about it on social media'],
         'correct': 1},
        {'question': f'Which of these is a common indicator of a {cat_display} attack?',
         'options': _quiz_indicators(cat),
         'correct': 0},
        {'question': f'What organizational measure best prevents {cat_display} incidents?',
         'options': _quiz_prevention(cat),
         'correct': 0},
        {'question': 'How often should cybersecurity awareness training be conducted?',
         'options': ['Once during onboarding only', 'Annually at minimum, with quarterly refreshers', 'Only after a security incident', 'Every five years'],
         'correct': 1},
    ]
    return questions

def _quiz_indicators(cat):
    indicators = {
        'phishing': ['Urgent language and mismatched sender domains', 'Slow internet connection', 'New software updates available', 'Calendar invitations from colleagues'],
        'ransomware': ['Files becoming inaccessible with ransom demands', 'Computer running slowly', 'Receiving spam email', 'Browser homepage changing'],
        'social_engineering': ['Unsolicited requests for sensitive information with urgency', 'Receiving marketing emails', 'Colleagues asking about projects', 'IT department sending newsletters'],
        'password_security': ['Multiple failed login attempts and account lockouts', 'Forgetting your password occasionally', 'Password expiration reminders', 'New password policy announcements'],
        'data_privacy': ['Unauthorized access logs and data exfiltration alerts', 'Cookie consent popups on websites', 'Privacy policy updates from services', 'GDPR compliance emails'],
        'zero_day_threats': ['Unusual system behavior with no known vulnerability patch available', 'Regular software update notifications', 'Antivirus scan completing successfully', 'Firewall blocking known threats'],
    }
    return indicators.get(cat, ['Unexpected system behavior', 'Normal operations', 'Routine maintenance', 'Scheduled updates'])

def _quiz_prevention(cat):
    prevention = {
        'phishing': ['Email filtering, SPF/DKIM/DMARC, and regular simulation exercises', 'Blocking all external email', 'Using only phone communication', 'Disabling email attachments'],
        'ransomware': ['Regular offline backups, network segmentation, and patch management', 'Paying ransoms quickly to minimize downtime', 'Using only cloud storage', 'Avoiding all file sharing'],
        'social_engineering': ['Security awareness training with verification procedures', 'Eliminating all phone communications', 'Restricting all visitor access permanently', 'Using only written communication'],
        'password_security': ['MFA enforcement, password managers, and credential monitoring', 'Requiring password changes every week', 'Writing passwords in a secure notebook', 'Using the same strong password everywhere'],
        'data_privacy': ['Data classification, encryption, and access controls with audit logging', 'Storing all data in a single location', 'Deleting all data after use', 'Making all data publicly accessible for transparency'],
        'zero_day_threats': ['Defense-in-depth with behavioral detection and rapid patching', 'Disconnecting from the internet', 'Only using open-source software', 'Waiting for vendors to release patches'],
    }
    return prevention.get(cat, ['Comprehensive security program', 'Single security tool', 'Annual audits only', 'Outsourcing all security'])

def _make_scenario(cat):
    scenario_pool = {
        'phishing': [
            {'situation':'You receive an email from your CEO requesting an urgent wire transfer of $50,000 to a new vendor. The email address looks correct but uses a slightly different domain (ceo@company-corp.com instead of ceo@companycorp.com).','choices':['Process the transfer immediately — it\'s from the CEO','Call the CEO directly on their known phone number to verify','Reply to the email asking for confirmation','Forward to your manager for approval'],'correct':1,'explanation':'Business Email Compromise (BEC) is a sophisticated phishing attack. Always verify unusual financial requests through a separate communication channel, never by replying to the suspicious email.'},
            {'situation':'You receive a text message claiming to be from your bank saying your account has been compromised. It includes a link to "verify your identity" and asks you to act within 2 hours.','choices':['Click the link and enter your credentials quickly','Call your bank using the number on your physical card, not the text','Reply STOP to block the sender','Forward the text to your colleagues as a warning'],'correct':1,'explanation':'Smishing (SMS phishing) creates artificial urgency. Never click links in unsolicited texts. Always contact your bank through their official app or the number on your card.'},
        ],
        'ransomware': [
            {'situation':'Multiple employees report that their files have been encrypted and a ransom note demands 5 Bitcoin within 48 hours. The IT team confirms the ransomware has spread to shared drives.','choices':['Negotiate with the attackers for a lower ransom','Immediately isolate affected systems, activate incident response plan, and contact law enforcement','Pay the ransom to get files back quickly','Attempt to decrypt the files using online tools'],'correct':1,'explanation':'The priority is containment. Isolate affected systems to stop lateral movement, activate your incident response plan, preserve evidence for forensics, and contact law enforcement. Paying ransoms funds criminal operations and doesn\'t guarantee recovery.'},
            {'situation':'An employee accidentally opens a macro-enabled Word document from an unknown sender. Their antivirus flags suspicious behavior but the employee dismisses the alert.','choices':['No action needed since antivirus caught it','Immediately disconnect the device from the network and report to IT security','Run a quick scan and continue working','Restart the computer to clear any issues'],'correct':1,'explanation':'Dismissed antivirus alerts can indicate partial malware execution. Immediately isolate the device and let the security team investigate. Early detection and isolation prevent ransomware from spreading across the network.'},
        ],
        'social_engineering': [
            {'situation':'A person in a delivery uniform arrives at your office lobby carrying packages. They ask an employee to hold the door open to the secure area, claiming they have a delivery for the 5th floor.','choices':['Hold the door — they look legitimate with the uniform and packages','Politely decline and direct them to reception to get a visitor badge','Ask to see their delivery order but still let them in','Call the 5th floor to ask if they\'re expecting a delivery and let them wait'],'correct':1,'explanation':'Tailgating is a common physical social engineering attack. Regardless of appearance, all visitors must go through proper check-in procedures. Direct them to reception where their identity and purpose can be verified.'},
            {'situation':'You receive a LinkedIn message from a recruiter at a prestigious company offering a role with 40% higher salary. They ask you to fill out an "application form" that requests your current employer\'s VPN credentials for a "background check."','choices':['Fill out the form — it\'s a great opportunity','Decline and report the profile as fraudulent','Provide only basic information, not VPN credentials','Share the opportunity with colleagues first'],'correct':1,'explanation':'Legitimate recruiters never ask for employer credentials. This is a targeted social engineering attack to gain access to your organization\'s network. Report the fraudulent profile immediately.'},
        ],
        'password_security': [
            {'situation':'During a routine security audit, it\'s discovered that 40% of employees are using passwords that appear in known breach databases. Several use patterns like CompanyName2024! or Season+Year combinations.','choices':['Send a company-wide email asking people to change passwords','Deploy a password manager, enforce MFA, and implement real-time breach monitoring','Add more complexity requirements (special characters, numbers)','Lock all affected accounts immediately without notice'],'correct':1,'explanation':'Simply requiring password changes leads to predictable patterns. A comprehensive approach includes password managers for unique passwords, MFA as a second defense layer, and continuous monitoring against known breach databases.'},
            {'situation':'A colleague shares their screen during a video call and their password manager is visible showing login credentials for a critical production system.','choices':['Ignore it — password managers are secure','Privately message the colleague to close the password manager and suggest they rotate the exposed credentials','Report the colleague to HR immediately','Take a screenshot for evidence'],'correct':1,'explanation':'Accidental credential exposure requires immediate but discreet action. Notify the colleague privately, have them rotate the exposed credentials, and review screen-sharing practices in security training.'},
        ],
        'data_privacy': [
            {'situation':'Your marketing team wants to use customer purchase data to train a machine learning model for personalized recommendations. The data includes names, emails, purchase history, and payment information.','choices':['Proceed — it\'s company data and can be used internally','Anonymize/pseudonymize the data, remove payment info, verify consent covers this use case, and conduct a privacy impact assessment','Only remove payment information and proceed','Ask customers for blanket consent through a terms update'],'correct':1,'explanation':'Data privacy regulations require purpose limitation and data minimization. Before repurposing data, verify consent covers the new use, anonymize where possible, remove unnecessary fields (especially payment data), and document the privacy impact.'},
            {'situation':'A former employee requests deletion of all their personal data from company systems under GDPR Article 17 (Right to Erasure). However, some of their data is in compliance-required audit logs.','choices':['Delete everything immediately to comply','Evaluate the request: delete personal data where legally permissible, retain only what\'s required by law with documentation, and respond within 30 days','Ignore the request — they\'re no longer an employee','Delete only their email account and consider it done'],'correct':1,'explanation':'The Right to Erasure has exceptions for legal obligations. Conduct a proper assessment, delete what you can, document legitimate reasons for retention (e.g., audit requirements), and respond to the data subject within the GDPR timeline.'},
        ],
        'zero_day_threats': [
            {'situation':'Your threat intelligence feed reports a critical zero-day vulnerability in a widely-used library that your web applications depend on. No vendor patch is available yet, and active exploitation has been confirmed in the wild.','choices':['Wait for the vendor to release an official patch','Implement compensating controls: WAF rules, network segmentation, enhanced monitoring, and consider disabling the affected functionality if possible','Take all affected systems offline immediately','Ignore it until exploitation is confirmed against your specific systems'],'correct':1,'explanation':'When no patch exists, implement defense-in-depth: deploy virtual patches via WAF, increase monitoring and logging, segment network access to affected systems, and evaluate whether the vulnerable functionality can be temporarily disabled. Coordinate with your vendor for patch timeline.'},
            {'situation':'A security researcher privately discloses a zero-day vulnerability in your organization\'s public-facing application, giving you 90 days before public disclosure.','choices':['Ignore it — researchers just want attention','Acknowledge the report, assign a team to develop and test a fix, coordinate disclosure timeline with the researcher, and deploy the patch before public disclosure','Threaten legal action against the researcher','Wait until the 90 days are almost up to begin working on a fix'],'correct':1,'explanation':'Responsible vulnerability disclosure should be taken seriously. Acknowledge promptly, assign engineering resources, develop and test a fix, deploy before public disclosure, and consider a bug bounty program. Coordinating with the researcher maintains a good security reputation.'},
        ],
    }
    pool = scenario_pool.get(cat, scenario_pool['phishing'])
    return random.choice(pool)
